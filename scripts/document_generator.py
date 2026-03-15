"""
Document Generator for Biz-Seed
Generates business-specific PowerPoint pitch decks and PDF financial statements.
"""

import os
from datetime import datetime
from typing import Dict, Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    print("⚠️  python-pptx not installed. Run: pip install python-pptx")
    PPTX_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
    from reportlab.lib.units import inch
    from reportlab.lib import colors as rl_colors
    from reportlab.lib.enums import TA_CENTER, TA_RIGHT
    PDF_AVAILABLE = True
except ImportError:
    print("⚠️  reportlab not installed. Run: pip install reportlab")
    PDF_AVAILABLE = False


# Brand colours
BLUE       = RGBColor(0,  102, 204)   # Primary
DARK_BLUE  = RGBColor(0,   70, 150)
GREEN      = RGBColor(34, 139,  34)   # Success / profit
WHITE      = RGBColor(255, 255, 255)
LIGHT_GREY = RGBColor(245, 245, 245)
DARK_TEXT  = RGBColor(30,  30,  30)
ORANGE     = RGBColor(255, 140,   0)  # Accent


BUSINESS_TYPE_LABELS = {
    'SPAZA_SHOP':      'Spaza Shop',
    'HOME_SALON':      'Home Salon',
    'STREET_VENDOR':   'Street Vendor',
    'HOME_BAKER':      'Home Bakery',
    'CAR_WASH':        'Car Wash',
    'TAILOR':          'Tailoring & Alterations',
    'COBBLER':         'Shoe Repair',
    'CATERING':        'Catering',
    'TAVERN':          'Tavern',
    'GENERAL':         'General Business',
}

BUSINESS_TYPE_PROBLEM = {
    'SPAZA_SHOP':    'Informal spaza shops serve millions of South Africans but are locked out of formal credit, bulk supplier deals, and business development support.',
    'HOME_SALON':    'Township hair stylists and beauty practitioners have loyal customers and real skills but no formal proof of income to access equipment loans or studio funding.',
    'STREET_VENDOR': 'Street vendors operate daily in high-foot-traffic areas with consistent sales but cannot access formal markets or working-capital finance.',
    'HOME_BAKER':    'Home bakers serve their communities but lack the capital for commercial equipment, formal kitchens, and retail shelf space.',
    'CAR_WASH':      'Township car-wash operators have steady customers but lack access to proper infrastructure, equipment finance, and business registration support.',
    'TAILOR':        'Skilled township tailors and seamstresses serve their communities but struggle to scale without access to sewing machines, fabric credit, or formal premises.',
    'GENERAL':       'Township micro-entrepreneurs generate consistent revenue but are excluded from formal financial services due to lack of documentation and credit history.',
}

BUSINESS_TYPE_SOLUTION = {
    'SPAZA_SHOP':    'We operate a community spaza shop providing essential groceries, airtime, and household goods to township residents at accessible prices, accepting both cash and digital payments.',
    'HOME_SALON':    'We provide hair and beauty services directly in the community, eliminating transport barriers and offering affordable treatments with flexible payment options.',
    'STREET_VENDOR': 'We provide convenient, affordable products at high-traffic community points, serving daily customer needs with both cash and digital payment acceptance.',
    'HOME_BAKER':    'We bake and supply fresh, affordable baked goods (bread, vetkoek, cakes) to our community, schools, and local events, using quality ingredients and community trust.',
    'CAR_WASH':      'We provide affordable, high-quality vehicle cleaning services to our township community, building regular client relationships through quality and convenience.',
    'TAILOR':        'We provide professional tailoring, alterations, and custom clothing services, combining skilled craftsmanship with community-level accessibility and pricing.',
    'GENERAL':       'We provide essential products and services to our township community, building customer loyalty through quality, trust, and community-focused operations.',
}


class DocumentGenerator:
    """Generate investor-ready documents specific to each business."""

    def __init__(self):
        self.output_dir = os.path.join('..', 'downloads', 'bizseed')
        os.makedirs(self.output_dir, exist_ok=True)

    # ─────────────────────────────────────────────────────────────────────────
    # INTERNAL HELPERS
    # ─────────────────────────────────────────────────────────────────────────

    def _add_slide(self, prs, bg_color=None):
        """Add a blank slide with optional background colour."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        if bg_color:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
        return slide

    def _textbox(self, slide, left, top, width, height, text, size, bold=False,
                 color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
        """Add a text box and return the paragraph for further styling."""
        txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        txb.text_frame.word_wrap = wrap
        tf = txb.text_frame
        tf.text = text
        p = tf.paragraphs[0]
        p.alignment = align
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        return p

    def _bullet_slide(self, prs, title_text, bullets, accent_color=BLUE):
        """Add a titled slide with bullet points."""
        slide = self._add_slide(prs, LIGHT_GREY)

        # Left accent bar
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0), Inches(0.15), Inches(7.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()

        # Title
        self._textbox(slide, 0.4, 0.3, 9.3, 0.8, title_text,
                      28, bold=True, color=accent_color)

        # Horizontal rule
        line = slide.shapes.add_shape(1, Inches(0.4), Inches(1.15), Inches(9.2), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = accent_color
        line.line.fill.background()

        # Bullets
        txb = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(9), Inches(5.5))
        txb.text_frame.word_wrap = True
        first = True
        for bullet in bullets:
            if first:
                p = txb.text_frame.paragraphs[0]
                first = False
            else:
                p = txb.text_frame.add_paragraph()
            p.text = f"  {bullet}"
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(8)

        return slide

    # ─────────────────────────────────────────────────────────────────────────
    # PITCH DECK
    # ─────────────────────────────────────────────────────────────────────────

    def generate_pitch_deck(self, business_data: Dict[str, Any]) -> str:
        """
        Generate a business-specific PowerPoint pitch deck.

        Required keys in business_data:
            businessName, businessType, location, ownerName,
            revenue, expenses, profit, transactionCount
        Optional:
            dailySales, empowerScore, userId, growthRate
        """
        if not PPTX_AVAILABLE:
            raise Exception("python-pptx not installed. Run: pip install python-pptx")

        prs = Presentation()
        prs.slide_width  = Inches(10)
        prs.slide_height = Inches(7.5)

        name         = business_data.get('businessName', 'Our Business')
        btype_key    = business_data.get('businessType', 'GENERAL')
        btype_label  = BUSINESS_TYPE_LABELS.get(btype_key, btype_key.replace('_', ' ').title())
        location     = business_data.get('location', 'South Africa')
        owner        = business_data.get('ownerName', 'The Owner')
        revenue      = float(business_data.get('revenue', 0))
        expenses     = float(business_data.get('expenses', 0))
        profit       = float(business_data.get('profit', 0))
        tx_count     = int(business_data.get('transactionCount', 0))
        daily_sales  = float(business_data.get('dailySales', revenue / 30 if revenue else 0))
        empower      = business_data.get('empowerScore', 'N/A')
        growth_rate  = float(business_data.get('growthRate', 10))
        margin       = (profit / revenue * 100) if revenue > 0 else 0
        ask_amount   = max(10000, round(revenue * 2 / 1000) * 1000)  # 2× monthly rev, min R10k

        problem_text  = BUSINESS_TYPE_PROBLEM.get(btype_key, BUSINESS_TYPE_PROBLEM['GENERAL'])
        solution_text = BUSINESS_TYPE_SOLUTION.get(btype_key, BUSINESS_TYPE_SOLUTION['GENERAL'])

        # ── Slide 1: Title ────────────────────────────────────────────────
        s1 = self._add_slide(prs, BLUE)

        # White diagonal accent shape (decorative)
        acc = s1.shapes.add_shape(1, Inches(6.5), Inches(0), Inches(5), Inches(7.5))
        acc.fill.solid()
        acc.fill.fore_color.rgb = DARK_BLUE
        acc.line.fill.background()

        self._textbox(s1, 0.6, 1.5, 6.5, 1.2, name,
                      40, bold=True, color=WHITE)
        self._textbox(s1, 0.6, 2.9, 6.5, 0.6, f"{btype_label}  •  {location}",
                      20, color=RGBColor(180, 220, 255))
        self._textbox(s1, 0.6, 3.8, 6.5, 0.5, "Investor Pitch Deck",
                      18, color=WHITE)
        self._textbox(s1, 0.6, 4.5, 6.5, 0.4, f"Presented by {owner}",
                      14, color=RGBColor(200, 230, 255))
        self._textbox(s1, 0.6, 6.8, 6.5, 0.4, datetime.now().strftime('%B %Y'),
                      12, color=RGBColor(180, 210, 255))

        # ── Slide 2: Executive Summary ────────────────────────────────────
        self._bullet_slide(prs, "Executive Summary", [
            f"{name} is a {btype_label} based in {location}",
            f"Owned and operated by {owner}",
            f"R{revenue:,.0f} in revenue over the last 30 days",
            f"R{profit:,.0f} net profit — {margin:.1f}% profit margin",
            f"{tx_count} customer transactions completed",
            f"MzansiPulse EmpowerScore: {empower} / 1000",
            "Seeking investment to formalise and scale operations",
        ], BLUE)

        # ── Slide 3: The Problem ──────────────────────────────────────────
        self._bullet_slide(prs, "The Problem", [
            problem_text,
            "",
            "Key barriers faced by township micro-entrepreneurs:",
            "  • No formal credit history or collateral",
            "  • No digital financial records accepted by banks",
            "  • Excluded from bulk purchasing agreements",
            "  • Limited access to business development support",
        ], ORANGE)

        # ── Slide 4: Our Solution ─────────────────────────────────────────
        self._bullet_slide(prs, "Our Solution", [
            solution_text,
            "",
            f"  • Serving the community of {location}",
            "  • Accepting cash and digital payments (inclusive)",
            "  • Tracked and verified via MzansiPulse platform",
            "  • Building formal financial history every transaction",
        ], GREEN)

        # ── Slide 5: Market Opportunity ───────────────────────────────────
        self._bullet_slide(prs, "Market Opportunity", [
            "South Africa's township economy is valued at R900 billion+",
            "Over 3 million micro-enterprises operate in informal settlements",
            "Less than 5% have access to formal business credit",
            "",
            f"Our immediate market: {location} community",
            f"  • {tx_count} transactions proves consistent local demand",
            f"  • Average daily sales of R{daily_sales:,.0f}",
            "  • High repeat customers = proven product-market fit",
        ], BLUE)

        # ── Slide 6: Traction ─────────────────────────────────────────────
        s6 = self._add_slide(prs, LIGHT_GREY)

        bar6 = s6.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
        bar6.fill.solid()
        bar6.fill.fore_color.rgb = GREEN
        bar6.line.fill.background()

        self._textbox(s6, 0.4, 0.3, 9.3, 0.8, "Traction & Metrics",
                      28, bold=True, color=GREEN)
        rule = s6.shapes.add_shape(1, Inches(0.4), Inches(1.15), Inches(9.2), Inches(0.04))
        rule.fill.solid()
        rule.fill.fore_color.rgb = GREEN
        rule.line.fill.background()

        # 4 metric cards
        metrics = [
            ("Monthly Revenue", f"R{revenue:,.0f}", BLUE),
            ("Monthly Profit",  f"R{profit:,.0f}",  GREEN),
            ("Transactions",    str(tx_count),        ORANGE),
            ("Profit Margin",   f"{margin:.1f}%",     DARK_BLUE),
        ]
        for i, (label, value, clr) in enumerate(metrics):
            col = i % 2
            row = i // 2
            left = 0.5 + col * 4.8
            top  = 1.6 + row * 2.4

            card = s6.shapes.add_shape(1, Inches(left), Inches(top), Inches(4.3), Inches(2.0))
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = clr
            card.line.width = Pt(2)

            txb = s6.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15),
                                         Inches(4.0), Inches(0.5))
            txb.text_frame.paragraphs[0].text = label
            txb.text_frame.paragraphs[0].font.size = Pt(13)
            txb.text_frame.paragraphs[0].font.color.rgb = DARK_TEXT

            txb2 = s6.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.65),
                                          Inches(4.0), Inches(0.9))
            txb2.text_frame.paragraphs[0].text = value
            txb2.text_frame.paragraphs[0].font.size = Pt(32)
            txb2.text_frame.paragraphs[0].font.bold = True
            txb2.text_frame.paragraphs[0].font.color.rgb = clr

        # ── Slide 7: 12-Month Financial Projection ────────────────────────
        proj_revenue  = revenue * ((1 + growth_rate / 100) ** 12)
        proj_profit   = proj_revenue * (margin / 100) if margin > 0 else proj_revenue * 0.25

        self._bullet_slide(prs, "12-Month Financial Projection", [
            f"Current monthly revenue:  R{revenue:,.0f}",
            f"Current monthly profit:   R{profit:,.0f}",
            f"Assumed growth rate:      {growth_rate:.0f}% per month (conservative)",
            "",
            f"Projected revenue in 12 months:  R{proj_revenue:,.0f}",
            f"Projected profit in 12 months:   R{proj_profit:,.0f}",
            "",
            "Growth drivers: stock expansion, digital marketing, vendor partnerships",
        ], DARK_BLUE)

        # ── Slide 8: The Ask ──────────────────────────────────────────────
        self._bullet_slide(prs, f"The Ask  —  R{ask_amount:,.0f}", [
            f"We are seeking R{ask_amount:,.0f} in funding to:",
            "",
            "  1. Expand stock/inventory to meet demand",
            "  2. Register business formally (CIPC + SARS)",
            "  3. Upgrade equipment and workspace",
            "  4. Build 3-month working-capital reserve",
            "",
            "Return on investment: breakeven within 6 months at current growth rate",
            f"Contact: {owner}  •  {location}",
        ], ORANGE)

        # ── Slide 9: Why Us ───────────────────────────────────────────────
        self._bullet_slide(prs, "Why Invest in Us", [
            f"Proven revenue: R{revenue:,.0f} earned last month with zero formal marketing",
            f"Consistent demand: {tx_count} completed transactions",
            "MzansiPulse verified: all transactions digitally recorded and auditable",
            "Community anchor: loyal local customer base with high repeat rate",
            f"Experienced operator: {owner} has hands-on knowledge of this market",
            "Fundable: EmpowerScore of {empower} demonstrates financial discipline".format(empower=empower),
            "Scalable: clear path from informal to formal business within 12 months",
        ], BLUE)

        # ── Slide 10: Contact ─────────────────────────────────────────────
        s10 = self._add_slide(prs, DARK_BLUE)
        self._textbox(s10, 1, 2.0, 8, 1.0, "Thank You",
                      44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        self._textbox(s10, 1, 3.2, 8, 0.6, name,
                      24, bold=True, color=RGBColor(180, 220, 255), align=PP_ALIGN.CENTER)
        self._textbox(s10, 1, 4.0, 8, 0.5, owner,
                      18, color=WHITE, align=PP_ALIGN.CENTER)
        self._textbox(s10, 1, 4.6, 8, 0.5, location,
                      16, color=RGBColor(160, 200, 240), align=PP_ALIGN.CENTER)
        self._textbox(s10, 1, 6.5, 8, 0.4,
                      "Generated by MzansiPulse  •  mzansipulse.app",
                      11, color=RGBColor(120, 160, 200), align=PP_ALIGN.CENTER)

        # Save
        filename = f"pitch_deck_{business_data.get('userId', 'demo')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(self.output_dir, filename)
        prs.save(filepath)
        print(f"✅ Pitch deck generated: {filepath}")
        return filename  # return just the filename so the API can build the URL

    # ─────────────────────────────────────────────────────────────────────────
    # FINANCIAL STATEMENTS (PDF)
    # ─────────────────────────────────────────────────────────────────────────

    def generate_financial_statements(self, business_data: Dict[str, Any]) -> str:
        if not PDF_AVAILABLE:
            raise Exception("reportlab not installed. Run: pip install reportlab")

        filename = f"financials_{business_data.get('userId', 'demo')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        filepath = os.path.join(self.output_dir, filename)

        doc = SimpleDocTemplate(filepath, pagesize=A4,
                                topMargin=0.75*inch, bottomMargin=0.75*inch)
        styles = getSampleStyleSheet()
        story  = []

        title_style = ParagraphStyle('Title', parent=styles['Title'],
                                     fontSize=22, spaceAfter=6)
        h2_style    = ParagraphStyle('H2', parent=styles['Heading2'],
                                     fontSize=14, textColor=rl_colors.HexColor('#0066CC'))
        normal      = styles['Normal']

        story.append(Paragraph(f"<b>{business_data['businessName']}</b>", title_style))
        story.append(Paragraph("Financial Statements — Last 30 Days", normal))
        story.append(Spacer(1, 0.3*inch))
        story.append(HRFlowable(width="100%", thickness=2,
                                color=rl_colors.HexColor('#0066CC')))
        story.append(Spacer(1, 0.2*inch))

        revenue  = float(business_data.get('revenue', 0))
        expenses = float(business_data.get('expenses', 0))
        profit   = revenue - expenses
        margin   = (profit / revenue * 100) if revenue > 0 else 0

        story.append(Paragraph("<b>Income Statement</b>", h2_style))
        income_data = [
            ['Description', 'Amount'],
            ['Total Revenue (Sales)', f"R {revenue:>10,.2f}"],
            ['Total Expenses (Cost of Goods + Operations)', f"R {expenses:>10,.2f}"],
            ['', ''],
            ['Net Profit / (Loss)', f"R {profit:>10,.2f}"],
            ['Profit Margin', f"{margin:.1f}%"],
        ]
        t = Table(income_data, colWidths=[4.5*inch, 2*inch])
        t.setStyle(TableStyle([
            ('BACKGROUND',   (0, 0), (-1, 0), rl_colors.HexColor('#0066CC')),
            ('TEXTCOLOR',    (0, 0), (-1, 0), rl_colors.white),
            ('FONTNAME',     (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE',     (0, 0), (-1, -1), 11),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1),
             [rl_colors.white, rl_colors.HexColor('#f5f5f5')]),
            ('LINEABOVE',    (0, 4), (-1, 4), 1.5, rl_colors.black),
            ('FONTNAME',     (0, 4), (-1, 4), 'Helvetica-Bold'),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
            ('TOPPADDING',   (0, 0), (-1, -1), 8),
        ]))
        story.append(t)
        story.append(Spacer(1, 0.4*inch))

        story.append(Paragraph("<b>Notes</b>", h2_style))
        story.append(Paragraph(
            f"These statements cover the 30-day period ending {datetime.now().strftime('%d %B %Y')}. "
            f"All figures are derived from transaction records captured on the MzansiPulse platform "
            f"and are provided for investor reference purposes.",
            normal
        ))

        doc.build(story)
        print(f"✅ Financial statements generated: {filepath}")
        return filename

    # ─────────────────────────────────────────────────────────────────────────
    # BUSINESS PLAN (PDF)
    # ─────────────────────────────────────────────────────────────────────────

    def generate_business_plan(self, business_data: Dict[str, Any]) -> str:
        """
        Generate a structured business plan PDF.
        Required keys: businessName, businessType, location, ownerName,
                       revenue, expenses, profit, transactionCount
        """
        if not PDF_AVAILABLE:
            raise Exception("reportlab not installed. Run: pip install reportlab")

        filename = f"business_plan_{business_data.get('userId', 'demo')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        filepath = os.path.join(self.output_dir, filename)

        doc = SimpleDocTemplate(filepath, pagesize=A4,
                                topMargin=0.75 * inch, bottomMargin=0.75 * inch,
                                leftMargin=inch, rightMargin=inch)

        styles  = getSampleStyleSheet()
        title_s = ParagraphStyle('BPTitle', parent=styles['Title'],
                                 fontSize=24, spaceAfter=6,
                                 textColor=rl_colors.HexColor('#0066CC'))
        h2_s    = ParagraphStyle('BPH2', parent=styles['Heading2'],
                                 fontSize=14, spaceBefore=18, spaceAfter=6,
                                 textColor=rl_colors.HexColor('#0066CC'))
        h3_s    = ParagraphStyle('BPH3', parent=styles['Heading3'],
                                 fontSize=12, spaceBefore=10, spaceAfter=4,
                                 textColor=rl_colors.HexColor('#004499'))
        body_s  = ParagraphStyle('BPBody', parent=styles['Normal'],
                                 fontSize=11, leading=16, spaceAfter=6)
        meta_s  = ParagraphStyle('BPMeta', parent=styles['Normal'],
                                 fontSize=10, textColor=rl_colors.grey)

        name      = business_data.get('businessName', 'Our Business')
        btype_key = business_data.get('businessType', 'GENERAL')
        btype     = BUSINESS_TYPE_LABELS.get(btype_key, btype_key.replace('_', ' ').title())
        location  = business_data.get('location', 'South Africa')
        owner     = business_data.get('ownerName', 'The Owner')
        revenue   = float(business_data.get('revenue', 0))
        expenses  = float(business_data.get('expenses', 0))
        profit    = revenue - expenses
        margin    = (profit / revenue * 100) if revenue > 0 else 0
        tx_count  = int(business_data.get('transactionCount', 0))
        ask       = max(10000, round(revenue * 2 / 1000) * 1000)

        problem  = BUSINESS_TYPE_PROBLEM.get(btype_key, BUSINESS_TYPE_PROBLEM['GENERAL'])
        solution = BUSINESS_TYPE_SOLUTION.get(btype_key, BUSINESS_TYPE_SOLUTION['GENERAL'])

        story = []

        # Cover
        story.append(Paragraph(f"<b>{name}</b>", title_s))
        story.append(Paragraph("Business Plan", styles['Heading2']))
        story.append(Spacer(1, 0.1 * inch))
        story.append(Paragraph(f"{btype}  •  {location}", meta_s))
        story.append(Paragraph(f"Prepared by: {owner}", meta_s))
        story.append(Paragraph(f"Date: {datetime.now().strftime('%B %Y')}", meta_s))
        story.append(Spacer(1, 0.2 * inch))
        story.append(HRFlowable(width="100%", thickness=3,
                                color=rl_colors.HexColor('#0066CC')))
        story.append(Spacer(1, 0.3 * inch))

        # 1. Executive Summary
        story.append(Paragraph("1. Executive Summary", h2_s))
        story.append(Paragraph(
            f"{name} is a {btype} operated by {owner} in {location}. "
            f"The business serves the local township community with consistent demand, "
            f"having recorded {tx_count} customer transactions. "
            f"Monthly revenue stands at <b>R{revenue:,.0f}</b> with a net profit of "
            f"<b>R{profit:,.0f}</b> ({margin:.1f}% margin). "
            f"This business plan outlines our strategy to formalise, grow, and access "
            f"investment funding of <b>R{ask:,.0f}</b>.",
            body_s))

        # 2. Business Description
        story.append(Paragraph("2. Business Description", h2_s))
        story.append(Paragraph(f"<b>Business Name:</b> {name}", body_s))
        story.append(Paragraph(f"<b>Business Type:</b> {btype}", body_s))
        story.append(Paragraph(f"<b>Owner / Operator:</b> {owner}", body_s))
        story.append(Paragraph(f"<b>Location:</b> {location}", body_s))
        story.append(Paragraph(
            solution,
            body_s))

        # 3. Problem & Opportunity
        story.append(Paragraph("3. Problem & Market Opportunity", h2_s))
        story.append(Paragraph(problem, body_s))
        story.append(Paragraph(
            "South Africa's township economy exceeds R900 billion annually. "
            "Over 3 million micro-enterprises operate informally with less than 5% "
            "having access to formal business credit. This gap represents a significant "
            "investment opportunity for funders willing to back proven, community-rooted businesses.",
            body_s))

        # 4. Products & Services
        story.append(Paragraph("4. Products & Services", h2_s))
        story.append(Paragraph(
            f"{name} offers {btype.lower()} services to the {location} community. "
            "We accept both cash and digital payments, making our services accessible "
            "to the full spectrum of community members. All transactions are recorded "
            "on the MzansiPulse platform, providing verified financial data.",
            body_s))

        # 5. Marketing Strategy
        story.append(Paragraph("5. Marketing Strategy", h2_s))
        story.append(Paragraph("<b>Current channels:</b>", h3_s))
        for point in [
            "Word-of-mouth referrals from existing loyal customers",
            "Physical presence in the community (high foot-traffic location)",
            "WhatsApp groups and community networks",
        ]:
            story.append(Paragraph(f"• {point}", body_s))
        story.append(Paragraph("<b>Growth plan with investment:</b>", h3_s))
        for point in [
            "Social media presence (Facebook, TikTok) for local awareness",
            "Loyalty programme to increase repeat business",
            "Partnership with local schools and community organisations",
            "Listing on local township business directories",
        ]:
            story.append(Paragraph(f"• {point}", body_s))

        # 6. Financial Summary
        story.append(Paragraph("6. Financial Summary", h2_s))
        fin_data = [
            ['Metric', 'Current (Monthly)', 'Projected (Month 12)'],
            ['Revenue',   f"R{revenue:,.0f}",  f"R{revenue * 2.5:,.0f}"],
            ['Expenses',  f"R{expenses:,.0f}", f"R{expenses * 2.0:,.0f}"],
            ['Net Profit', f"R{profit:,.0f}",  f"R{profit * 3.0:,.0f}"],
            ['Profit Margin', f"{margin:.1f}%", f"{min(margin + 5, 45):.1f}%"],
            ['Transactions', str(tx_count), f"{tx_count * 3}+"],
        ]
        tbl = Table(fin_data, colWidths=[2.2 * inch, 2 * inch, 2 * inch])
        tbl.setStyle(TableStyle([
            ('BACKGROUND',   (0, 0), (-1, 0), rl_colors.HexColor('#0066CC')),
            ('TEXTCOLOR',    (0, 0), (-1, 0), rl_colors.white),
            ('FONTNAME',     (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE',     (0, 0), (-1, -1), 10),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1),
             [rl_colors.white, rl_colors.HexColor('#f0f4ff')]),
            ('GRID',         (0, 0), (-1, -1), 0.5, rl_colors.HexColor('#cccccc')),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 7),
            ('TOPPADDING',   (0, 0), (-1, -1), 7),
        ]))
        story.append(tbl)

        # 7. Funding Requirements
        story.append(Paragraph("7. Funding Requirements", h2_s))
        story.append(Paragraph(
            f"<b>{name} is seeking R{ask:,.0f} in funding</b> to accelerate growth "
            f"and formalise operations. The funds will be used as follows:",
            body_s))
        use_of_funds = [
            ['Use of Funds', 'Allocation'],
            ['Stock / Inventory Expansion',     '35%'],
            ['CIPC Registration & Legal Costs',  '10%'],
            ['Equipment Upgrade',               '25%'],
            ['Working Capital Reserve (3 months)', '20%'],
            ['Marketing & Digital Presence',    '10%'],
        ]
        tbl2 = Table(use_of_funds, colWidths=[4 * inch, 1.5 * inch])
        tbl2.setStyle(TableStyle([
            ('BACKGROUND',   (0, 0), (-1, 0), rl_colors.HexColor('#0066CC')),
            ('TEXTCOLOR',    (0, 0), (-1, 0), rl_colors.white),
            ('FONTNAME',     (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE',     (0, 0), (-1, -1), 10),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1),
             [rl_colors.white, rl_colors.HexColor('#f0f4ff')]),
            ('GRID',         (0, 0), (-1, -1), 0.5, rl_colors.HexColor('#cccccc')),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 7),
            ('TOPPADDING',   (0, 0), (-1, -1), 7),
        ]))
        story.append(tbl2)
        story.append(Spacer(1, 0.3 * inch))
        story.append(Paragraph(
            f"<i>Generated by MzansiPulse on {datetime.now().strftime('%d %B %Y')}. "
            f"All financial figures are derived from verified platform transaction data.</i>",
            meta_s))

        doc.build(story)
        print(f"✅ Business plan generated: {filepath}")
        return filename

    # ─────────────────────────────────────────────────────────────────────────
    # GROWTH FORECAST (kept as before — returns filename)
    # ─────────────────────────────────────────────────────────────────────────

    def generate_growth_forecast(self, business_data: Dict[str, Any]) -> str:
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment
            from openpyxl.chart import LineChart, Reference
        except ImportError:
            raise Exception("openpyxl not installed. Run: pip install openpyxl")

        filename = f"forecast_{business_data.get('userId', 'demo')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        filepath = os.path.join(self.output_dir, filename)

        wb = Workbook()
        ws = wb.active
        ws.title = "12-Month Forecast"

        header_fill = PatternFill(start_color="0066CC", end_color="0066CC", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True)

        for col, heading in enumerate(['Month', 'Revenue (R)', 'Expenses (R)', 'Profit (R)'], 1):
            cell = ws.cell(row=1, column=col, value=heading)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal='center')

        current_revenue = float(business_data.get('revenue', 0))
        expense_ratio   = 0.65
        growth_rate     = 1.08

        for month in range(1, 13):
            rev  = round(current_revenue * (growth_rate ** (month - 1)), 2)
            exp  = round(rev * expense_ratio, 2)
            prof = round(rev - exp, 2)
            ws.cell(row=month + 1, column=1, value=f"Month {month}")
            ws.cell(row=month + 1, column=2, value=rev)
            ws.cell(row=month + 1, column=3, value=exp)
            ws.cell(row=month + 1, column=4, value=prof)

        chart = LineChart()
        chart.title = f"12-Month Forecast — {business_data.get('businessName', 'Business')}"
        chart.y_axis.title = "Amount (R)"
        chart.x_axis.title = "Month"
        chart.add_data(Reference(ws, min_col=2, min_row=1, max_row=13, max_col=4),
                       titles_from_data=True)
        chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=13))
        ws.add_chart(chart, "F2")

        wb.save(filepath)
        print(f"✅ Growth forecast generated: {filepath}")
        return filename
